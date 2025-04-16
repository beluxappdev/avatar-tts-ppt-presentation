from flask import Flask, request, render_template, jsonify, send_file, redirect, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
import os
import uuid
import json
import logging
import time
from presentation_generator import PresentationGenerator

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='[%(asctime)s] [%(levelname)s] %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Configuration
UPLOAD_FOLDER = 'uploads'
PROCESSED_FOLDER = 'processed'
ALLOWED_EXTENSIONS = {'pptx', 'ppt'}

# Create necessary directories
for folder in [UPLOAD_FOLDER, PROCESSED_FOLDER]:
    os.makedirs(folder, exist_ok=True)

# Dictionary to store active presentation generators
active_generators = {}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        # Generate a unique identifier for this upload
        process_id = str(uuid.uuid4())
        
        # Create a directory for this process
        process_dir = os.path.join(PROCESSED_FOLDER, process_id)
        os.makedirs(process_dir, exist_ok=True)
        
        # Save the uploaded file
        filename = secure_filename(file.filename)
        filepath = os.path.join(UPLOAD_FOLDER, f"{process_id}_{filename}")
        file.save(filepath)
        
        # Process the PowerPoint file and extract slide notes
        try:
            slide_notes = extract_slide_notes(filepath)
            
            # Save the slide-notes data to a JSON file
            data_file = os.path.join(process_dir, 'slide_notes.json')
            with open(data_file, 'w') as f:
                json.dump(slide_notes, f, indent=2)
            
            return jsonify({
                'success': True,
                'process_id': process_id,
                'slides_count': len(slide_notes),
                'slides': slide_notes
            })
            
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            return jsonify({'error': f"Error processing file: {str(e)}"}), 500
    
    return jsonify({'error': 'Invalid file format'}), 400

@app.route('/generate-video/<process_id>', methods=['POST'])
def generate_video(process_id):
    """Generate a complete presentation video with avatars reading slide notes"""
    logger.info(f"Received video generation request for process: {process_id}")
    
    # Get the process directory
    process_dir = os.path.join(PROCESSED_FOLDER, process_id)
    if not os.path.exists(process_dir):
        logger.warning(f"Process directory not found: {process_dir}")
        return jsonify({'error': 'Process not found'}), 404
    
    # Load the slide notes data
    data_file = os.path.join(process_dir, 'slide_notes.json')
    if not os.path.exists(data_file):
        logger.warning(f"Slide notes file not found: {data_file}")
        return jsonify({'error': 'Slide notes not found'}), 404
    
    with open(data_file, 'r') as f:
        slide_notes = json.load(f)
        logger.info(f"Loaded slide notes, found {len(slide_notes)} slides")
    
    # Find the original PowerPoint file
    pptx_files = [f for f in os.listdir(UPLOAD_FOLDER) if f.startswith(process_id) and f.endswith(('.pptx', '.ppt'))]
    if not pptx_files:
        logger.warning(f"No PowerPoint file found for process ID: {process_id}")
        return jsonify({'error': 'PowerPoint file not found'}), 404
    
    pptx_file_path = os.path.join(UPLOAD_FOLDER, pptx_files[0])
    logger.info(f"Using PowerPoint file: {pptx_file_path}")
    
    # Get avatar settings from the request
    avatar_settings = {
        "voice": request.form.get('voice', 'en-US-JennyNeural'),
        "character": request.form.get('character', 'lisa'),
        "style": request.form.get('style', 'graceful-sitting')
    }
    
    # Get max concurrency parameter
    try:
        max_concurrency = int(request.form.get('max_concurrency', '2'))
        if max_concurrency < 1:
            max_concurrency = 1
    except ValueError:
        max_concurrency = 2
        
    logger.info(f"Starting presentation generation with max concurrency: {max_concurrency}")
    
    # Create a presentation generator
    generator = PresentationGenerator(
        process_id=process_id,
        pptx_file_path=pptx_file_path,
        output_dir=process_dir,
        max_concurrency=max_concurrency
    )
    active_generators[process_id] = generator
    
    # Start generating the presentation
    result = generator.generate_presentation(slide_notes, avatar_settings)
    return jsonify(result)

@app.route('/status/<process_id>', methods=['GET'])
def get_status(process_id):
    """Get the status of a presentation generation process"""
    logger.info(f"Checking status for process: {process_id}")
    
    # Check if we have an active generator for this process
    if process_id in active_generators:
        status = active_generators[process_id].get_status()
        return jsonify(status)
    
    # Try to load status from file
    status_file = os.path.join(PROCESSED_FOLDER, process_id, 'status.json')
    if os.path.exists(status_file):
        try:
            with open(status_file, 'r') as f:
                status = json.load(f)
            return jsonify(status)
        except Exception as e:
            logger.error(f"Error reading status file: {str(e)}")
            return jsonify({'error': f"Error reading status: {str(e)}"}), 500
    
    return jsonify({'error': 'Process not found'}), 404

@app.route('/download/<process_id>', methods=['GET'])
def download_video(process_id):
    """Download the generated presentation video from Azure Blob Storage"""
    logger.info(f"Received download request for process: {process_id}")
    
    # First check if we have a local file
    output_file = os.path.join(PROCESSED_FOLDER, f"{process_id}_presentation.mp4")
    
    if os.path.exists(output_file):
        try:
            return send_file(
                output_file,
                mimetype='video/mp4',
                as_attachment=True,
                download_name=f"presentation_{process_id}.mp4"
            )
        except Exception as e:
            logger.error(f"Error downloading local presentation video: {str(e)}")
    
    # If no local file, check for Azure blob URL in status file
    status_file = os.path.join(PROCESSED_FOLDER, process_id, 'status.json')
    if os.path.exists(status_file):
        try:
            with open(status_file, 'r') as f:
                status = json.load(f)
                
            # Check if concatenation status exists and contains output URL
            if ('jobs' in status and 
                'concatenation_status' in status['jobs'] and 
                status['jobs']['concatenation_status'].get('status') == 'completed'):
                
                # Check if we have a blob URL
                if 'blob_url' in status['jobs']['concatenation_status']:
                    blob_url = status['jobs']['concatenation_status']['blob_url']
                    logger.info(f"Redirecting to blob storage URL: {blob_url}")
                    return redirect(blob_url)
        except Exception as e:
            logger.error(f"Error reading status file: {str(e)}")
    
    return jsonify({'error': 'Presentation video file not found'}), 404

def extract_slide_notes(pptx_file_path):
    """Extract notes from all slides in a PowerPoint file"""
    logger.info(f"Extracting notes from PowerPoint file: {pptx_file_path}")
    
    slide_notes = []
    presentation = Presentation(pptx_file_path)
    
    for i, slide in enumerate(presentation.slides):
        # Extract notes
        notes = ""
        if hasattr(slide, 'notes_slide') and slide.notes_slide is not None:
            notes = slide.notes_slide.notes_text_frame.text
        
        # Add to the list
        slide_notes.append({
            "slide_index": i + 1,
            "notes": notes
        })
        
        logger.info(f"Processed slide {i+1} with {len(notes)} characters of notes")
    
    return slide_notes

if __name__ == '__main__':
    # Run on all interfaces so it's accessible from outside the container
    app.run(host='0.0.0.0', port=5000, debug=False)