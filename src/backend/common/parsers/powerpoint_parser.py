import os
import io
import tempfile
import subprocess
import logging
from typing import List, Dict
import pdf2image # type: ignore
from PIL import Image # type: ignore
from pptx import Presentation # type: ignore

logger = logging.getLogger(__name__)


class PowerPointParser:
    """Parser for PowerPoint files - extracting images, notes, and other content"""
    
    def __init__(self, dpi: int = 150):
        """Initialize the PowerPoint parser
        
        Args:
            dpi (int): DPI for image conversion (default: 150 for good quality vs file size balance)
        """
        self.dpi = dpi
    
    def convert_to_images(self, ppt_file_data: bytes) -> List[Image.Image]:
        """Convert PowerPoint file data to a list of PIL Images
        
        Args:
            ppt_file_data (bytes): The PowerPoint file as bytes
            
        Returns:
            List[Image.Image]: List of PIL Image objects, one per slide
            
        Raises:
            Exception: If conversion fails at any step
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save the PowerPoint file
            ppt_path = os.path.join(temp_dir, "presentation.pptx")
            with open(ppt_path, "wb") as f:
                f.write(ppt_file_data)
            
            # Convert to PDF first
            pdf_path = self._convert_ppt_to_pdf(ppt_path, temp_dir)
            if not pdf_path:
                raise Exception("Failed to convert PowerPoint to PDF")
            
            # Convert PDF to images
            return self._convert_pdf_to_images(pdf_path)
    
    def extract_notes(self, ppt_data: bytes) -> List[Dict[str, any]]:
        """Extract notes from PowerPoint file data
        
        Args:
            ppt_data (bytes): The PowerPoint file as bytes
            
        Returns:
            List of dictionaries containing slide index and notes text
        """
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_path = temp_file.name
            temp_file.write(ppt_data)
        
        try:
            # Load the presentation using python-pptx
            presentation = Presentation(temp_path)
            
            # Extract notes from each slide
            slide_notes = []
            
            for slide_index, slide in enumerate(presentation.slides):
                notes_text = ""
                has_notes = False
                
                try:
                    # Check if slide has notes
                    if hasattr(slide, 'notes_slide') and slide.notes_slide:
                        if hasattr(slide.notes_slide, 'notes_placeholder') and slide.notes_slide.notes_placeholder:
                            notes_text = slide.notes_slide.notes_placeholder.text
                            has_notes = bool(notes_text.strip())  # Only set true if there's actual text content
                    
                    if not has_notes:
                        notes_text = "This is a default script"
                        has_notes = True 

                    slide_info = {
                        "index": slide_index,
                        "notes_text": notes_text.strip() if notes_text else "",
                        "has_notes": has_notes
                    }
                    
                    slide_notes.append(slide_info)
                    
                    if has_notes:
                        logger.info(f"Found notes for slide {slide_index}: {len(notes_text)} characters")
                    else:
                        logger.info(f"No notes found for slide {slide_index}")
                        
                except Exception as e:
                    logger.error(f"Error extracting notes from slide {slide_index}: {str(e)}")
                    # Add slide with error state
                    slide_notes.append({
                        "index": slide_index,
                        "notes_text": "",
                        "has_notes": False,
                        "error": str(e)
                    })
            
            logger.info(f"Extracted notes from {len(slide_notes)} slides")
            return slide_notes
            
        except Exception as e:
            logger.error(f"Error loading PowerPoint presentation: {str(e)}")
            raise Exception(f"Failed to extract notes from PowerPoint: {str(e)}")
        finally:
            # Clean up temporary file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def get_slide_count(self, ppt_data: bytes) -> int:
        """Get the number of slides in a PowerPoint presentation
        
        Args:
            ppt_data (bytes): The PowerPoint file as bytes
            
        Returns:
            int: Number of slides in the presentation
        """
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_path = temp_file.name
            temp_file.write(ppt_data)
        
        try:
            presentation = Presentation(temp_path)
            return len(presentation.slides)
        except Exception as e:
            logger.error(f"Error getting slide count: {str(e)}")
            raise Exception(f"Failed to get slide count from PowerPoint: {str(e)}")
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def image_to_bytes(self, image: Image.Image, format: str = 'PNG') -> bytes:
        """Convert PIL Image to bytes
        
        Args:
            image (Image.Image): PIL Image object
            format (str): Image format (PNG, JPEG, etc.)
            
        Returns:
            bytes: Image data as bytes
        """
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format=format)
        img_byte_arr.seek(0)
        return img_byte_arr.getvalue()
    
    def _convert_ppt_to_pdf(self, ppt_path: str, output_dir: str) -> str:
        """Convert PowerPoint to PDF using LibreOffice
        
        Args:
            ppt_path (str): Path to PowerPoint file
            output_dir (str): Directory to save the output PDF
            
        Returns:
            str: Path to the output PDF file
            
        Raises:
            Exception: If conversion fails
        """
        logger.info(f"Converting PowerPoint to PDF: {ppt_path}")
        
        try:
            # Run LibreOffice headless to convert PowerPoint to PDF
            command = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                ppt_path
            ]
            
            # Execute the command
            process = subprocess.run(
                command,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=True,
                timeout=300  # 5 minute timeout
            )
            
            logger.info(f"LibreOffice conversion completed")
            if process.stdout:
                logger.debug(f"LibreOffice output: {process.stdout}")
            
            # Get the output PDF path
            ppt_filename = os.path.basename(ppt_path)
            pdf_filename = os.path.splitext(ppt_filename)[0] + '.pdf'
            pdf_path = os.path.join(output_dir, pdf_filename)
            
            # Check if the file exists
            if os.path.exists(pdf_path):
                logger.info(f"Successfully converted PowerPoint to PDF: {pdf_path}")
                return pdf_path
            else:
                raise Exception(f"PDF file not created: {pdf_path}")
                
        except subprocess.CalledProcessError as e:
            error_msg = f"LibreOffice conversion failed: {e.stderr}"
            logger.error(error_msg)
            raise Exception(error_msg)
        except subprocess.TimeoutExpired:
            error_msg = "LibreOffice conversion timed out"
            logger.error(error_msg)
            raise Exception(error_msg)
        except Exception as e:
            error_msg = f"Unexpected error in PowerPoint to PDF conversion: {str(e)}"
            logger.error(error_msg)
            raise Exception(error_msg)
    
    def _convert_pdf_to_images(self, pdf_path: str) -> List[Image.Image]:
        """Convert PDF to images
        
        Args:
            pdf_path (str): Path to PDF file
            
        Returns:
            List[Image.Image]: List of PIL Image objects
            
        Raises:
            Exception: If conversion fails
        """
        logger.info(f"Converting PDF to images: {pdf_path}")
        
        try:
            # Convert PDF to images using pdf2image
            images = pdf2image.convert_from_path(
                pdf_path,
                dpi=self.dpi,
                fmt='png'
            )
            
            logger.info(f"Successfully converted PDF to {len(images)} images")
            return images
            
        except Exception as e:
            error_msg = f"Error converting PDF to images: {str(e)}"
            logger.error(error_msg)
            raise Exception(error_msg)