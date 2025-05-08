import logging
from datetime import datetime

from pptx import Presentation
from base_extractor import BaseExtractor

# Configure logging
logger = logging.getLogger(__name__)

class ScriptExtractor(BaseExtractor):
    """extractor for extracting scripts (notes) from PowerPoint presentations"""
    
    def __init__(self):
        """Initialize the script extractor"""
        super().__init__(extractor_type="script")

    def _process_powerpoint(self, ppt_id, file_path):
        """Process PowerPoint file to extract scripts
        
        Args:
            ppt_id (str): ID of the PowerPoint
            file_path (str): Path to the downloaded PowerPoint file
        """
        # Load the presentation using python-pptx
        presentation = Presentation(file_path)
        
        # Extract scripts from the PowerPoint
        self.extract_scripts(ppt_id, presentation)
    
    def extract_scripts(self, ppt_id, presentation):
        """Extract scripts (notes) from PowerPoint presentation
        
        Args:
            ppt_id (str): ID of the PowerPoint
            presentation (Presentation): PowerPoint presentation object
        """
        logger.info(f"Extracting scripts from PowerPoint: {ppt_id}")
        
        # Track results for each slide for error reporting
        slide_results = {}
        overall_success = True
        error_messages = []
        
        try:
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_result = {
                    "index": slide_index,
                    "success": True,
                    "has_notes": False,
                    "errors": []
                }
                
                logger.info(f"Processing slide {slide_index + 1} of {len(presentation.slides)}")
                
                try:
                    # Create the slide directory path
                    slide_dir_path = f"{ppt_id}/slides/{slide_index}"
                    
                    # Extract and save notes if available
                    notes_text = ""
                    has_notes = False
                    
                    # Check if slide has notes
                    if hasattr(slide, 'notes_slide') and slide.notes_slide:
                        logger.info(f"Slide {slide_index} has notes shape: {slide.notes_slide.shapes}")
                        logger.info(f"Slide {slide_index} has notes placeholder: {slide.notes_slide.notes_placeholder}")
                        if hasattr(slide.notes_slide, 'notes_placeholder') and slide.notes_slide.notes_placeholder:
                            notes_text = slide.notes_slide.notes_placeholder.text
                            has_notes = bool(notes_text.strip())  # Only set true if there's actual text content
                    
                    # Save notes to blob storage if they exist
                    if has_notes:
                        notes_blob_path = f"{slide_dir_path}/script.txt"
                        notes_blob_client = self.blob_container_client.get_blob_client(notes_blob_path)
                        
                        try:
                            notes_blob_client.upload_blob(notes_text.encode('utf-8'), overwrite=True)
                            logger.info(f"Uploaded notes: {notes_blob_path}")
                            slide_result["has_notes"] = True
                        except Exception as e:
                            error_msg = f"Error uploading notes for slide {slide_index}: {str(e)}"
                            logger.error(error_msg)
                            slide_result["errors"].append(error_msg)
                            slide_result["success"] = False
                    else:
                        logger.info(f"No notes found for slide {slide_index}")
                    
                except Exception as e:
                    error_msg = f"Error processing notes for slide {slide_index}: {str(e)}"
                    logger.error(error_msg)
                    slide_result["success"] = False
                    slide_result["errors"].append(error_msg)
                    overall_success = False
                
                # Store the result for this slide
                slide_results[slide_index] = slide_result
            
            # Store additional metadata about slide extraction
            self.extraction_results = {
                "ppt_id": ppt_id,
                "overall_success": overall_success,
                "slide_count": len(presentation.slides),
                "slides_processed": len(slide_results),
                "slides_with_errors": sum(1 for result in slide_results.values() if not result["success"]),
                "slides_with_notes": sum(1 for result in slide_results.values() if result.get("has_notes", False)),
                "slide_results": slide_results,
                "error_messages": error_messages
            }
            
            logger.info(f"Completed script extraction from PowerPoint: {ppt_id} - Success: {overall_success}")
            
        except Exception as e:
            error_msg = f"Critical error extracting scripts: {str(e)}"
            logger.error(error_msg)
            error_messages.append(error_msg)
            overall_success = False
            
            # Store failure metadata
            self.extraction_results = {
                "ppt_id": ppt_id,
                "overall_success": False,
                "error_messages": error_messages
            }
            
            raise

    def _update_extraction_data(self, item, ppt_id):
        """Update script extraction data in the Cosmos DB document

        Args:
            item (dict): Cosmos DB document to update
            ppt_id (str): ID of the PowerPoint
        """
        try:
            slide_prefix = f"{ppt_id}/slides/"
            slides_data = item.get('slides', []) or []
        
            # Create a dictionary of existing slides indexed by slide number
            slides_dict = {slide.get('index'): slide for slide in slides_data if 'index' in slide}
        
            # Dictionary to track which slides have scripts
            slides_with_scripts = {}
        
            # List all blobs with the prefix for this presentation (done once)
            blobs = list(self.blob_container_client.list_blobs(name_starts_with=slide_prefix))
        
            # Process all blobs in a single pass
            for blob in blobs:
                # Get relative path after the slide prefix
                rel_path = blob.name[len(slide_prefix):]
                parts = rel_path.split('/')
            
                if len(parts) < 1:
                    continue
                
                try:
                    slide_index = int(parts[0])
                
                    # Check if this is a script file
                    if len(parts) > 1 and parts[1] == "script.txt":
                        script_blob_path = blob.name
                    
                        # Create or retrieve slide data
                        if slide_index not in slides_dict:
                            slides_dict[slide_index] = {
                                "index": slide_index,
                                "hasImage": False,
                                "hasScript": True
                            }
                        else:
                            slides_dict[slide_index]["hasScript"] = True
                    
                        # Add script information
                        slides_dict[slide_index]["scriptUrl"] = f"https://{self.blob_endpoint.replace('https://', '')}/{self.blob_container_name}/{script_blob_path}"
                        slides_dict[slide_index]["scriptSize"] = blob.size
                        slides_with_scripts[slide_index] = True
                    
                        # Fetch script content for preview if it's not too large
                        if blob.size < 10000:  # Only fetch if less than 10KB
                            script_blob_client = self.blob_container_client.get_blob_client(script_blob_path)
                            script_content = script_blob_client.download_blob().readall().decode('utf-8')
                            slides_dict[slide_index]["scriptContent"] = script_content
                
                    # Process other slide-related data here if needed
                    # (e.g., check for images or other slide components)
                
                except ValueError:
                    logger.warning(f"Invalid slide directory format: {parts[0]}")
        
            # Convert the dictionary back to a sorted list
            processed_slides = [slides_dict[index] for index in sorted(slides_dict.keys())]
        
            # Update the document with processed slide data
            item['slides'] = processed_slides
            item['slideCount'] = len(processed_slides)
            item['slidesWithNotes'] = len(slides_with_scripts)

        except Exception as e:
            # Log the error but don't fail the entire operation
            logger.error(f"Error updating script data: {str(e)}")

            # Update the document with error information
            item['scriptProcessingStatus'] = "Failed" 
            item['scriptProcessingError'] = str(e)
            item['scriptProcessingErrorAt'] = datetime.utcnow().isoformat()

            raise    

def main():
    """Main entry point"""
    try:
        extractor = ScriptExtractor()
        extractor.start_processing()
    except Exception as e:
        logger.error(f"Unhandled exception in script extractor: {str(e)}")
        raise

if __name__ == "__main__":
    main()